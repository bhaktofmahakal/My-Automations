import os
from pptx import Presentation

ops_pptx = r"u:\My-Automations\vocallabs-assignment\Vocallabs_Founders_Office_Utsav\2_Presentation_Operational_Answers\vocallabs_ops_presentation.pptx"
inv_pptx = r"u:\My-Automations\vocallabs-assignment\Vocallabs_Founders_Office_Utsav\3_Investment_Pitch_Deck\vocallabs_investment_deck.pptx"

def inspect_deck(filename, deck_name):
    prs = Presentation(filename)
    print(f"\n==========================================")
    print(f"   FULL TEXT INSPECTION: {deck_name}")
    print(f"==========================================")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- SLIDE {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        print(f"  [P] {paragraph.text.strip()}")

inspect_deck(ops_pptx, "OPERATIONAL ANSWERS DECK (Q1, Q2, Q3)")
inspect_deck(inv_pptx, "INVESTMENT PITCH DECK")
