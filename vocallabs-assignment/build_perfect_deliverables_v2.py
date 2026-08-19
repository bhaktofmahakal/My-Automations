import sys
import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

output_dir = r"u:\My-Automations\vocallabs-assignment\02_pitch_deck"
drive_dir = r"u:\My-Automations\vocallabs-assignment\Vocallabs_Founders_Office_Utsav"
visuals_dir = r"u:\My-Automations\vocallabs-assignment\03_visuals"

# Color Palette (Executive Dark Theme)
C_DARK_BG   = RGBColor(11, 15, 23)     # #0B0F17 Deep Slate
C_CARD_BG   = RGBColor(26, 36, 56)    # #1A2438 Slate Card Fill
C_INDIGO    = RGBColor(129, 140, 248) # #818CF8 Soft Indigo
C_EMERALD   = RGBColor(52, 211, 153)  # #34D399 Bright Emerald
C_AMBER     = RGBColor(251, 191, 36)  # #FBBF24 Amber Gold
C_CORAL     = RGBColor(248, 113, 113) # #F87171 Coral Red
C_TEXT_MAIN = RGBColor(248, 250, 252)# #F8FAFC White Text
C_TEXT_MUTED= RGBColor(148, 163, 184)# #94A3B8 Muted Grey

def add_slide_header(slide, title, category="VOCALLABS AI | FOUNDER'S OFFICE STRATEGY"):
    tb1 = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = category.upper()
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.font.color.rgb = C_INDIGO
    
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_MAIN

def add_glass_card(slide, left, top, width, height, title, items, accent_color=C_INDIGO):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD_BG
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = accent_color
    p0.space_after = Pt(6)
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "- " + item
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_TEXT_MAIN
        p.space_after = Pt(4)

def add_kpi_card(slide, left, top, width, height, stat_num, stat_label, subtext, color=C_EMERALD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = stat_num
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = color
    
    p2 = tf.add_paragraph()
    p2.text = stat_label
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_MAIN
    
    p3 = tf.add_paragraph()
    p3.text = subtext
    p3.font.size = Pt(9)
    p3.font.color.rgb = C_TEXT_MUTED

# =========================================================
# DECK 1: OPERATIONAL STRATEGY PRESENTATION
# =========================================================
prs1 = Presentation()
prs1.slide_width = Inches(13.333)
prs1.slide_height = Inches(7.5)
b1 = prs1.slide_layouts[6]

# Slide 1: Cover
s1 = prs1.slides.add_slide(b1)
s1.background.fill.solid()
s1.background.fill.fore_color.rgb = C_DARK_BG

tb = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "VOCALLABS AI | FOUNDER'S OFFICE INTERN"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_INDIGO
p.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "Operational Execution Playbook & Strategic Answers"
p2.font.size = Pt(30)
p2.font.bold = True
p2.font.color.rgb = C_TEXT_MAIN
p2.space_after = Pt(12)

p3 = tf.add_paragraph()
p3.text = "Direct Answers to Q1 (Client Escalation), Q2 (INR 1 Cr Budget & Time), and Q3 (Candidate Positioning)"
p3.font.size = Pt(15)
p3.font.color.rgb = C_TEXT_MUTED
p3.space_after = Pt(24)

add_kpi_card(s1, Inches(0.8), Inches(4.8), Inches(2.7), Inches(1.8), "<400ms", "Audio Pipeline Latency", "STT -> LLM -> TTS stream", C_INDIGO)
add_kpi_card(s1, Inches(3.8), Inches(4.8), Inches(2.7), Inches(1.8), "INR 1.0 Cr", "Capital Deployment", "12-month growth budget", C_EMERALD)
add_kpi_card(s1, Inches(6.8), Inches(4.8), Inches(2.7), Inches(1.8), "40/30/20/10", "Capacity Prioritization", "Deals, Eng, CS & Ops split", C_AMBER)
add_kpi_card(s1, Inches(9.8), Inches(4.8), Inches(2.7), Inches(1.8), "10 / 10", "Candidate Alignment", "Ships code, decks & strategy", C_CORAL)

# Slide 2: Q1 Sub-Question 1 & 2 + 48h Visual Timeline Chart
s2 = prs1.slides.add_slide(b1)
s2.background.fill.solid()
s2.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(s2, "Q1: What I Would Do First & Who I Would Loop In Internally")

chart_path_timeline = os.path.join(visuals_dir, "escalation_timeline_chart.png")
if os.path.exists(chart_path_timeline):
    s2.shapes.add_picture(chart_path_timeline, Inches(0.8), Inches(1.5), width=Inches(5.7))

add_glass_card(s2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.6),
               "1. WHAT I WOULD DO FIRST (Hours 00 - 02 Triage)",
               ["Audit VocalStack Live Telemetry: Inspect past 30 days logs for latency spikes (>450ms), STT Word Error Rate (WER), and SIP packet loss.",
                "120-Minute Client SLA Message: Send direct memo to Client VP of Ops promising full root-cause breakdown by 4:00 PM.",
                "Isolate Audio Bottleneck: Pinpoint whether issue is LLM tool-calling lag, regional accent dropouts, or carrier packet jitter."],
               C_INDIGO)

add_glass_card(s2, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6),
               "2. WHO I WOULD LOOP IN INTERNALLY (Hours 02 - 12)",
               ["Solutions Engineer & AI Voice Architect: To analyze prompt context windows, STT acoustic models, and API callback latency.",
                "Telephony & Infrastructure Lead: To inspect WebRTC/SIP trunking routing, carrier failovers, and edge node performance.",
                "Co-Founders & CEO (5-Min Briefing): Brief Mritunjoy, Rag, & Nileesh to align on 30-day performance credit boundary."],
               C_EMERALD)

# Slide 3: Q1 Sub-Question 3
s3 = prs1.slides.add_slide(b1)
s3.background.fill.solid()
s3.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(s3, "Q1: How I Would Approach the Actual Client Conversation")

add_glass_card(s3, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "3. CLIENT CONVERSATION APPROACH (Minutes 00 - 25)",
               ["Radical Transparency (Mins 0-10): Present transparent telemetry logs showing where latency spiked: 100% ownership, zero blaming third-party APIs.",
                "Live Proof Demo (Mins 10-25): Demonstrate live side-by-side audio comparison of old pipeline vs newly optimized sub-400ms VocalStack pipeline."],
               C_AMBER)

add_glass_card(s3, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "3. CLIENT SUMMIT LOCK-IN (Minutes 25 - 50)",
               ["30-Day Performance Guarantee (Mins 25-40): Offer 50% credit on infra fees if metrics miss 25% target in 14 days + weekly joint engineering steering sync.",
                "Commercial Renewal (Mins 40-50): Convert crisis turnaround into long-term enterprise contract extension."],
               C_EMERALD)

# Slide 4: Q2 Budget
s4 = prs1.slides.add_slide(b1)
s4.background.fill.solid()
s4.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(s4, "Q2: INR 1 Crore Capital Deployment Strategy (Line-Item Breakdown)")

chart_path1 = os.path.join(visuals_dir, "budget_allocation_chart.png")
if os.path.exists(chart_path1):
    s4.shapes.add_picture(chart_path1, Inches(0.8), Inches(1.5), width=Inches(5.7))

add_glass_card(s4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "Capital Line-Item Allocation (INR 1,00,00,000 Total)",
               ["INR 40L (40%) Partner Acquisition Engine: Outbound B2B data enrichment for BPO targeting & funding first 10k partner voice call minutes.",
                "INR 35L (35%) Developer & OS Ecosystem: Contributor grants for VocalFlow OS & PocoDisk, production MCP Servers, and Voice AI hackathons.",
                "INR 15L (15%) Global Telephony Infra: Multi-region edge nodes (Mumbai, US East, EU) & multi-carrier SIP trunk failover.",
                "INR 10L (10%) Founder's Office Reserve: Key account turnaround credits & rapid micro-experimentation."],
               C_EMERALD)

# Slide 5: Q2 Time Split
s5 = prs1.slides.add_slide(b1)
s5.background.fill.solid()
s5.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(s5, "Q2: Founder's Office Weekly Time Prioritization Ratio")

chart_path2 = os.path.join(visuals_dir, "time_allocation_chart.png")
if os.path.exists(chart_path2):
    s5.shapes.add_picture(chart_path2, Inches(0.8), Inches(1.5), width=Inches(5.8))

add_glass_card(s5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "Weekly Operating Prioritization Ratio",
               ["40% (~18 hrs/wk) Partner Deals & GTM: Running outbound partner campaigns, pitching BPO owners, negotiating reseller margin splits.",
                "30% (~14 hrs/wk) Product & Eng Sync: Triaging partner feature requests, testing VocalStack/MCP releases, dogfooding Hiringg.",
                "20% (~9 hrs/wk) Partner Success & CS: Monitoring voice telemetry across top 20 accounts, enforcing zero churn.",
                "10% (~4 hrs/wk) Executive Operations: Metric dashboards, co-founder briefings, internal process automation."],
               C_INDIGO)

# Slide 6: Q3 Candidate Pitch
s6 = prs1.slides.add_slide(b1)
s6.background.fill.solid()
s6.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(s6, "Q3: Candidate Positioning & 10/10 Assessment Pitch")

add_glass_card(s6, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "Why I am a 10/10 Match for Vocallabs",
               ["Aligned with Mritunjoy's Ethos: I ship working code, PPTX engines, and operational frameworks: no 3-month slide deck fluff.",
                "Technical & Business Dual-Threat: Capable of analyzing voice telemetry (STT/TTS latency) AND negotiating gross margin splits with BPO founders.",
                "Product Meta-Signal Awareness: Recognized immediately that Hiringg (app.hiringg.ai) is Vocallabs' dogfooded AI hiring product.",
                "Zero Hand-Holding Needed: Self-directed operator who builds data-driven solutions under tight constraints."],
               C_EMERALD)

add_glass_card(s6, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "Operating Comparison: Me vs Standard Applicant",
               ["Standard Applicant: Writes 3 paragraphs or generic deck; lacks understanding of white-label infrastructure economics.",
                "My Execution: Built 2 complete PPTX presentations, strategy docs, visual telemetry charts, and teleprompter video script.",
                "Understanding of Moat: Focuses on VocalStack, VocalFlow OS, MCP Server, and BPO margin arbitrage.",
                "Final Pitch: An execution partner who runs beside the founders at 100mph."],
               C_AMBER)

ops_path = os.path.join(output_dir, "vocallabs_ops_presentation.pptx")
prs1.save(ops_path)

# Deck 2 Investment
prs2 = Presentation()
prs2.slide_width = Inches(13.333)
prs2.slide_height = Inches(7.5)
b2 = prs2.slide_layouts[6]

si1 = prs2.slides.add_slide(b2)
si1.background.fill.solid()
si1.background.fill.fore_color.rgb = C_DARK_BG
tb = si1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "INTERNAL STRATEGIC PROPOSAL & INVESTMENT DECK"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_INDIGO
p.space_after = Pt(10)
p2 = tf.add_paragraph()
p2.text = "Vocallabs: Scaling White-Label Voice AI Infrastructure"
p2.font.size = Pt(30)
p2.font.bold = True
p2.font.color.rgb = C_TEXT_MAIN
p2.space_after = Pt(12)
p3 = tf.add_paragraph()
p3.text = "Empowering BPOs, Agencies, and SaaS Platforms to Own & Resell Sovereign Voice AI Agents"
p3.font.size = Pt(15)
p3.font.color.rgb = C_TEXT_MUTED
p3.space_after = Pt(24)

add_kpi_card(si1, Inches(0.8), Inches(4.8), Inches(2.7), Inches(1.8), "INR 1.20 / min", "Wholesale Infra Rate", "Partner retails at INR 3-5/min", C_INDIGO)
add_kpi_card(si1, Inches(3.8), Inches(4.8), Inches(2.7), Inches(1.8), "62.5%", "Infra Gross Margin", "High operational leverage", C_EMERALD)
add_kpi_card(si1, Inches(6.8), Inches(4.8), Inches(2.7), Inches(1.8), "100+", "Active BPO Partners", "12-month partner target", C_AMBER)
add_kpi_card(si1, Inches(9.8), Inches(4.8), Inches(2.7), Inches(1.8), "INR 5.2 Cr", "Target Net ARR", "50M voice mins / month", C_CORAL)

si2 = prs2.slides.add_slide(b2)
si2.background.fill.solid()
si2.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(si2, "The Market Opportunity: The $150B Voice BPO Shift to Sovereign AI", "INVESTMENT DECK | MARKET OPPORTUNITY")
add_glass_card(si2, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "The Foil: Generic API Wrappers (Retell / Vapi)",
               ["Disintermediated Partners: Agencies & BPOs cannot build defensible businesses wrapping raw APIs: they get disintermediated.",
                "Loss of Brand Sovereignty: End-clients demand white-label solutions where the partner owns the customer relationship.",
                "High Latency & Jitter: Standard wrappers suffer from >1000ms latency and audio drops on international SIP routes.",
                "Destructive Pricing: High per-minute fees destroy partner gross margins on large-scale collections and sales drives."],
               C_AMBER)
add_glass_card(si2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "The Vocallabs Paradigm: White-Label Infrastructure",
               ["100% White-Label Sovereignty: Partners rebrand VocalStack as their own proprietary Voice AI platform.",
                "Sub-400ms Audio Pipeline: Integrated STT + Vocalassist AI reasoning + ultra-fast TTS streaming.",
                "Partner Margin Arbitrage: Partners buy infrastructure at wholesale rates (INR 1.20/min), retail at INR 3.00-INR 5.00/min, keeping 65% gross margins.",
                "Open Source Developer Moat: VocalFlow (OS visual flow builder) & PocoDisk (audio cache) drive viral developer adoption."],
               C_EMERALD)

si3 = prs2.slides.add_slide(b2)
si3.background.fill.solid()
si3.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(si3, "Product Architecture: The Integrated Vocallabs Engine", "INVESTMENT DECK | PRODUCT SUITE")
add_glass_card(si3, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
               "Core Voice Stack",
               ["VocalStack Voice: Sub-400ms ultra-low latency speech-to-speech runtime engine.",
                "Vocalassist AI: Context-aware reasoning, interruption handling & tool calling.",
                "Vocallabs Identity: Enterprise voice auth, biometric security & compliance."],
               C_INDIGO)
add_glass_card(si3, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
               "Developer & Open Source",
               ["MCP Server: Native Model Context Protocol for Claude & Cursor integrations.",
                "VocalFlow (OS): Visual flow builder for complex conversation workflows.",
                "PocoDisk (OS): High-throughput low-latency audio caching layer."],
               C_EMERALD)
add_glass_card(si3, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
               "Telephony & Dogfooding",
               ["Vocal Dialer: Multi-line high-concurrency predictive outbound dialer.",
                "Hiringg (app.hiringg.ai): Live dogfooded AI candidate interviewing platform.",
                "Enterprise Admin Portal: Multi-tenant white-label partner dashboard."],
               C_AMBER)

si4 = prs2.slides.add_slide(b2)
si4.background.fill.solid()
si4.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(si4, "Go-To-Market Engine: The Partner Reseller Flywheel", "INVESTMENT DECK | GO-TO-MARKET")
add_glass_card(si4, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "Target Partner Segments",
               ["Digital Marketing Agencies: Adding Voice AI sales agents to lead gen packages ($2k-$5k/mo retainer).",
                "Tier-2 BPOs & Contact Centers: Transitioning legacy call center seats to autonomous voice agents with zero CapEx.",
                "Vertical SaaS Platforms: Embedding native Voice AI capabilities into CRM, Healthcare, and Real Estate software.",
                "Collections & Debt Recovery: Deploying high-throughput automated payment reminder agents."],
               C_INDIGO)
add_glass_card(si4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "The B2B Growth Flywheel",
               ["Developer Virality: Open-source VocalFlow OS drives grassroots engineer discovery.",
                "Partner Onboarding: White-label sandbox credits (funded by INR 1 Cr budget) eliminate onboarding friction.",
                "Volume Monetization: As partners scale end-clients, Vocallabs infra minute consumption explodes.",
                "Defensible Lock-In: Deep API integrations & white-label admin dashboards ensure near-zero partner churn."],
               C_EMERALD)

si5 = prs2.slides.add_slide(b2)
si5.background.fill.solid()
si5.background.fill.fore_color.rgb = C_DARK_BG
add_slide_header(si5, "Financial Model & Unit Economics (Scaling to INR 5.2 Cr ARR)", "INVESTMENT DECK | UNIT ECONOMICS")
add_glass_card(si5, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "Partner Economics & Gross Margin Structure",
               ["Partner Wholesale Rate: INR 1.20 / minute (Vocallabs infra cost: ~INR 0.45 / min).",
                "Partner Resale Rate: INR 3.00 - INR 5.00 / minute (Partner retains ~65% gross margin).",
                "Vocallabs Infrastructure Gross Margin: 62.5% gross margin on runtime traffic.",
                "Average Partner Volume: 100,000 mins/mo = INR 1,20,000/mo net infra revenue per active partner."],
               C_EMERALD)
add_glass_card(si5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4),
               "12-Month ARR Target & Capital Deployment",
               ["Month 3 Target: 25 Active Partners | 10M Mins/mo | INR 1.2 Cr ARR.",
                "Month 6 Target: 50 Active Partners | 25M Mins/mo | INR 2.8 Cr ARR.",
                "Month 12 Target: 100 Active Partners | 50M Mins/mo | INR 5.2 Cr ARR.",
                "Capital Efficiency: Reaching profitability by Month 8 using INR 1 Cr strategic growth budget."],
               C_AMBER)

ops_path = os.path.join(output_dir, "vocallabs_ops_presentation.pptx")
prs1.save(ops_path)

inv_path = os.path.join(output_dir, "vocallabs_investment_deck.pptx")
prs2.save(inv_path)

# Sync to drive
shutil.copy(ops_path, os.path.join(drive_dir, "2_Presentation_Operational_Answers", "vocallabs_ops_presentation.pptx"))
shutil.copy(inv_path, os.path.join(drive_dir, "3_Investment_Pitch_Deck", "vocallabs_investment_deck.pptx"))

print("[SLIDE 2 VISUAL TIMELINE CHART ADDED & RE-GENERATED SUCCESS]")
